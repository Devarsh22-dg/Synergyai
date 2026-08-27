import os
import io
import re
import base64
import socket
import ipaddress
import zipfile
from urllib.parse import urlparse, urljoin
import requests
from bs4 import BeautifulSoup
import streamlit as st
import pandas as pd
import anthropic
import pypdf
from docx import Document
from pptx import Presentation
from PIL import Image
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill
from datetime import datetime

# auth is imported lazily inside check_access() rather than here — it pulls in
# streamlit_authenticator, which needs a real Streamlit component runtime.
# Importing it eagerly at module level would break headless imports of this
# file (e.g. the eval harness in evals/, which never calls check_access()).

# --- Configuration & Setup ---
st.set_page_config(layout="wide", page_title="ScopeForge: Consulting Accelerator")

DEFAULT_MODEL = "claude-sonnet-4-6"
FAST_MODEL = "claude-haiku-4-5-20251001"

# Display-only labels for the sidebar picker — the stored value stays the raw
# model ID, so which model actually gets called is unchanged.
MODEL_LABELS = {
    DEFAULT_MODEL: "Sonnet — best quality",
    FAST_MODEL: "Haiku — faster & cheaper",
}

STARTER_PROJECTS = ["Alpha-FinTech Migration", "Beta-Supply Chain Optimization", "Gamma-HR Platform Rollout"]
PROJECT_STATUSES = ["Planning", "In Progress", "On Hold", "Complete"]

DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"

# --- Access Control ---
# Built-in username/password, session-cookie login (see auth.py). No email
# verification step — accounts activate immediately on signup. MFA is a
# deliberately separate, later phase — not implemented here.
AUTH_ENABLED = True


def check_access():
    """Gate for the whole app. Renders the login/signup UI and returns False
    if nobody's logged in yet — callers should st.stop() in that case."""
    if not AUTH_ENABLED:
        return True
    import auth  # see note above the (removed) top-level import
    return auth.require_login()

CHATBOT_SYSTEM_PROMPT = (
    "You are ScopeBot, an AI assistant embedded in a tool used by business analysts, "
    "project managers, and program managers. Answer questions about requirements engineering, "
    "BRD/FRD best practices, Agile story writing, stakeholder management, and JIRA/Azure DevOps "
    "workflows. Keep answers practical and concise (a few short paragraphs or a brief list). "
    "If a question doesn't relate to those domains, answer briefly and steer back to how "
    "ScopeForge's modules (Elicitation Analysis, Documentation Generator, Story Creator, "
    "Meeting Actionizer) might help."
)

DOC_TYPE_CODES = {
    "BRD (Business Requirements Document)": "BRD",
    "FRD (Functional Requirements Document)": "FRD",
    "Data Dictionary": "Data_Dictionary",
    "Use Cases": "Use_Cases",
    "As-Is / To-Be Process Document": "AsIs_ToBe",
}

# --- Visual Theme ---
# Palette lives in theme.py so auth.py's login page uses the exact same values
# — see that module's docstring for why it's shared rather than duplicated.
from theme import (  # noqa: E402
    NAVY, NAVY_SOFT, ACCENT, ACCENT_HOVER, TEXT, TEXT_MUTED, TEXT_MUTED_STRONG,
    BORDER, SIDEBAR_BG, SIDEBAR_TEXT, SIDEBAR_MUTED,
)


def inject_theme():
    st.markdown(
        f"""
        <style>
        [data-testid="stHeader"] {{
            background-color: {SIDEBAR_BG};
            color: {SIDEBAR_TEXT};
        }}
        /* The header sits on the navy surface, but Streamlit still colors its
           controls for a light header: the sidebar expand chevron computes to
           rgba(17,24,39,.6) — near-black on near-black — so it vanishes
           completely. Setting `color` on the header alone does not reach it,
           because the icon span carries its own explicit color. */
        [data-testid="stHeader"] [data-testid="stIconMaterial"],
        [data-testid="stHeader"] svg {{
            color: {SIDEBAR_TEXT} !important;
            fill: {SIDEBAR_TEXT} !important;
        }}
        /* Contrast alone still left it reading as part of the bar — a bare
           glyph on a wide navy expanse has no affordance. Give the sidebar
           toggle an actual surface so it reads as a control. */
        [data-testid="stExpandSidebarButton"],
        [data-testid="stSidebarCollapseButton"] button,
        button[data-testid="stSidebarCollapseButton"] {{
            background-color: rgba(255, 255, 255, 0.10) !important;
            border: 1px solid rgba(255, 255, 255, 0.24) !important;
            border-radius: 6px !important;
            /* Streamlit hides the collapse chevron until the sidebar is
               hovered, so there is nothing to tell a first-time user the
               sidebar can be collapsed at all. Keep it on screen. */
            visibility: visible !important;
            opacity: 1 !important;
        }}
        [data-testid="stExpandSidebarButton"]:hover,
        [data-testid="stSidebarCollapseButton"] button:hover,
        button[data-testid="stSidebarCollapseButton"]:hover {{
            background-color: rgba(255, 255, 255, 0.20) !important;
            border-color: rgba(255, 255, 255, 0.42) !important;
        }}
        [data-testid="stAppViewContainer"] .main .block-container {{
            padding-top: 1.2rem !important;
        }}
        [data-testid="stSidebar"] {{
            background-color: {SIDEBAR_BG};
        }}
        [data-testid="stSidebar"] label,
        [data-testid="stSidebar"] p,
        [data-testid="stSidebar"] span,
        [data-testid="stSidebar"] div {{
            color: {SIDEBAR_TEXT} !important;
        }}
        [data-testid="stSidebar"] .stSelectbox div[data-baseweb="select"] > div,
        [data-testid="stSidebar"] .stRadio div[role="radiogroup"] {{
            background-color: {NAVY_SOFT};
        }}
        [data-testid="stSidebar"] hr {{
            border-color: {NAVY_SOFT};
        }}
        /* Six tabs overflow the container at ordinary window widths: the last
           one ("Traceability & Change Impact") clips off entirely, so a whole
           module is undiscoverable. Wrap the row rather than hiding tabs. */
        [data-testid="stTabs"] div[role="tablist"] {{
            flex-wrap: wrap;
            row-gap: 4px;
        }}
        [data-testid="stTabs"] button[role="tab"] {{
            background-color: #F3F4F6;
            color: {TEXT_MUTED_STRONG};
            border-radius: 8px 8px 0 0;
            padding: 0.55rem 1.1rem;
            font-weight: 600;
            margin-right: 4px;
        }}
        [data-testid="stTabs"] button[role="tab"] p {{
            color: inherit !important;
            font-weight: 600;
        }}
        [data-testid="stTabs"] button[aria-selected="true"] {{
            background-color: #FFFFFF;
            border-bottom: 3px solid {ACCENT};
        }}
        [data-testid="stTabs"] button[aria-selected="true"] p {{
            color: {ACCENT} !important;
        }}
        div[role="radiogroup"] {{
            gap: 0.4rem;
        }}
        /* Streamlit's default (non-"primary"-typed) buttons render in a plain
           dark/outline style. Nearly every button in this app is a primary
           action (Generate, Analyze, Sign Up, Download, ...) with no real
           secondary/de-emphasized button anywhere, so giving them all the one
           accent color reads as consistent rather than picking favorites. */
        button[data-testid^="stBaseButton-secondary"] {{
            background-color: {ACCENT};
            border: 1px solid {ACCENT};
            border-radius: 6px;
            font-weight: 600;
        }}
        /* Color the button element itself, not just a descendant <p>. The file
           uploader's "Browse files" puts its label in a bare text node with no
           <p> to hook, so a p-only rule left it at #111827 on the accent blue
           — 2.68:1, unreadable — while every other button looked correct. */
        button[data-testid^="stBaseButton-secondary"],
        button[data-testid^="stBaseButton-secondary"] p,
        button[data-testid^="stBaseButton-secondary"] span,
        button[data-testid^="stBaseButton-secondary"] div {{
            color: #FFFFFF !important;
        }}
        button[data-testid^="stBaseButton-secondary"]:hover {{
            background-color: {ACCENT_HOVER};
            border-color: {ACCENT_HOVER};
        }}
        </style>
        """,
        unsafe_allow_html=True,
    )


def render_masthead():
    """A branded top banner so the space at the top of the page reads as a
    deliberate header instead of dead whitespace."""
    st.markdown(
        f"""
        <div style="background-color: {SIDEBAR_BG}; padding: 0.9rem 1.5rem; border-radius: 10px;
                    margin-bottom: 1.3rem; display: flex; align-items: baseline; gap: 0.7rem;">
            <span style="font-size: 1.6rem; font-weight: 700; color: {SIDEBAR_TEXT}; letter-spacing: 0.2px;">ScopeForge</span>
            <span style="font-size: 0.9rem; color: {SIDEBAR_MUTED};">Consulting Accelerator for Business Analysts</span>
        </div>
        """,
        unsafe_allow_html=True,
    )


def section_header(title, subtitle):
    st.markdown(
        f"""
        <div style="border-left: 3px solid {ACCENT}; padding: 0.35rem 0 0.35rem 0.9rem; margin-bottom: 0.8rem;">
            <div style="font-size: 1.3rem; font-weight: 700; color: {TEXT}; line-height: 1.2;">{title}</div>
            <div style="font-size: 0.9rem; color: {TEXT_MUTED}; margin-top: 2px;">{subtitle}</div>
        </div>
        """,
        unsafe_allow_html=True,
    )


GAP_ANALYSIS_SCHEMA = {
    "type": "object",
    "properties": {
        "risk_score": {
            "type": "integer",
            "minimum": 0,
            "maximum": 100,
            "description": "Overall requirements risk score, 0 (low risk) to 100 (high risk).",
        },
        "risk_level": {"type": "string", "enum": ["Low", "Medium", "High", "Critical"]},
        "summary": {
            "type": "string",
            "description": "One or two sentence summary of the overall state of these requirements.",
        },
        "open_questions": {
            "type": "array",
            "items": {
                "type": "object",
                "properties": {
                    "type": {
                        "type": "string",
                        "enum": [
                            "Ambiguity", "Missing NFR", "Conflict",
                            "Missing Stakeholder Input", "Scope Risk", "Other",
                        ],
                    },
                    "issue": {"type": "string", "description": "Specific description of the gap or ambiguity."},
                    "why_it_matters": {"type": "string", "description": "Why this needs to be resolved."},
                },
                "required": ["type", "issue", "why_it_matters"],
            },
        },
    },
    "required": ["risk_score", "risk_level", "open_questions"],
}

STORY_SCHEMA = {
    "type": "object",
    "properties": {
        "stories": {
            "type": "array",
            "items": {
                "type": "object",
                "properties": {
                    "requirement": {"type": "string", "description": "Short label for the source requirement/need."},
                    "user_story": {
                        "type": "string",
                        "description": "Format: As a <role>, I want <capability>, so that <benefit>.",
                    },
                    "acceptance_criteria": {
                        "type": "string",
                        "description": "Gherkin-style GIVEN/WHEN/THEN acceptance criteria.",
                    },
                },
                "required": ["requirement", "user_story", "acceptance_criteria"],
            },
        }
    },
    "required": ["stories"],
}

MEETING_SCHEMA = {
    "type": "object",
    "properties": {
        "summary": {"type": "string", "description": "A short executive summary of the meeting."},
        "decisions": {"type": "array", "items": {"type": "string"}, "description": "Key decisions made."},
        "action_items": {
            "type": "array",
            "items": {
                "type": "object",
                "properties": {
                    "action": {"type": "string"},
                    "owner": {"type": "string", "description": "Use 'Unassigned' if no owner is stated."},
                    "due_date": {"type": "string", "description": "Use 'Not specified' if no date is stated."},
                },
                "required": ["action", "owner", "due_date"],
            },
        },
    },
    "required": ["summary", "decisions", "action_items"],
}

DATA_DICT_SCHEMA = {
    "type": "object",
    "properties": {
        "fields": {
            "type": "array",
            "items": {
                "type": "object",
                "properties": {
                    "field_name": {"type": "string"},
                    "data_type": {"type": "string"},
                    "description": {"type": "string"},
                    "source_system": {"type": "string", "description": "Use 'Not specified' if unknown."},
                    "validation_rules": {"type": "string", "description": "Use 'None specified' if unknown."},
                },
                "required": ["field_name", "data_type", "description", "source_system", "validation_rules"],
            },
        }
    },
    "required": ["fields"],
}

ASIS_TOBE_SCHEMA = {
    "type": "object",
    "properties": {
        "steps": {
            "type": "array",
            "items": {
                "type": "object",
                "properties": {
                    "step_id": {"type": "string", "description": "Sequential plain-number ID with no leading zeros, e.g. '1', '2', '3' (not '010')."},
                    "process_step": {"type": "string", "description": "Short name of this process step."},
                    "as_is_description": {"type": "string"},
                    "to_be_description": {"type": "string"},
                    "gap_or_change": {"type": "string", "description": "What needs to change to get from As-Is to To-Be."},
                    "shape_type": {"type": "string", "enum": ["Start", "Process", "Decision", "End"]},
                    "next_step_id": {
                        "type": "string",
                        "description": "Comma-separated step_id(s) this flows into next, no spaces, no leading zeros (e.g. '2' or '2,3'). Blank for the End step.",
                    },
                },
                "required": ["step_id", "process_step", "as_is_description", "to_be_description", "gap_or_change", "shape_type", "next_step_id"],
            },
        }
    },
    "required": ["steps"],
}

IMAGE_ANALYSIS_SCHEMA = {
    "type": "object",
    "properties": {
        "images": {
            "type": "array",
            "items": {
                "type": "object",
                "properties": {
                    "image_number": {"type": "integer", "description": "1-based index matching the order images were provided in."},
                    "content_type": {
                        "type": "string",
                        "enum": ["Screenshot/UI", "Diagram/Flowchart", "Chart/Graph", "Table", "Photo", "Other"],
                    },
                    "description": {
                        "type": "string",
                        "description": "What the image shows — for screenshots, describe the UI/screen and any labels, fields, or text visible; for diagrams, describe the flow/relationships; for charts, describe the data trend and axis labels.",
                    },
                },
                "required": ["image_number", "content_type", "description"],
            },
        }
    },
    "required": ["images"],
}

WORKSHOP_PREP_SCHEMA = {
    "type": "object",
    "properties": {
        "objectives": {"type": "string", "description": "1-2 sentences on what this workshop should achieve."},
        "agenda": {
            "type": "array",
            "items": {
                "type": "object",
                "properties": {
                    "topic": {"type": "string"},
                    "duration_minutes": {"type": "integer"},
                    "purpose": {"type": "string", "description": "Why this agenda item matters / what it should produce."},
                },
                "required": ["topic", "duration_minutes", "purpose"],
            },
        },
        "questions": {
            "type": "array",
            "items": {
                "type": "object",
                "properties": {
                    "category": {"type": "string", "description": "e.g. Scope, Process, Data, Constraints, Success Criteria."},
                    "question": {"type": "string"},
                },
                "required": ["category", "question"],
            },
        },
    },
    "required": ["objectives", "agenda", "questions"],
}

TEST_CASE_SCHEMA = {
    "type": "object",
    "properties": {
        "test_cases": {
            "type": "array",
            "items": {
                "type": "object",
                "properties": {
                    "test_id": {"type": "string", "description": "Sequential ID like TC-1, TC-2."},
                    "related_story": {"type": "string", "description": "Short label for the source user story."},
                    "scenario": {"type": "string", "description": "Short name of what this test verifies."},
                    "preconditions": {"type": "string", "description": "Use 'None' if not applicable."},
                    "steps": {"type": "string", "description": "Numbered test steps as a single string, e.g. '1. ... 2. ...'."},
                    "expected_result": {"type": "string"},
                    "priority": {"type": "string", "enum": ["High", "Medium", "Low"]},
                },
                "required": ["test_id", "related_story", "scenario", "preconditions", "steps", "expected_result", "priority"],
            },
        }
    },
    "required": ["test_cases"],
}

GLOSSARY_SCHEMA = {
    "type": "object",
    "properties": {
        "terms": {
            "type": "array",
            "items": {
                "type": "object",
                "properties": {
                    "term": {"type": "string"},
                    "definition": {"type": "string", "description": "Plain-language definition grounded in how the term is used in the source content."},
                    "source_context": {"type": "string", "description": "Brief note on where/how this term appears in the source. Use 'General' if not tied to a specific spot."},
                },
                "required": ["term", "definition", "source_context"],
            },
        }
    },
    "required": ["terms"],
}

PRIORITIZATION_SCHEMA = {
    "type": "object",
    "properties": {
        "items": {
            "type": "array",
            "items": {
                "type": "object",
                "properties": {
                    "requirement": {"type": "string", "description": "Short label for the requirement/need."},
                    "moscow_category": {"type": "string", "enum": ["Must Have", "Should Have", "Could Have", "Won't Have"]},
                    "rationale": {"type": "string", "description": "Why this requirement falls into this category."},
                },
                "required": ["requirement", "moscow_category", "rationale"],
            },
        }
    },
    "required": ["items"],
}

CHANGE_IMPACT_SCHEMA = {
    "type": "object",
    "properties": {
        "summary": {"type": "string", "description": "One-paragraph plain-language summary of the change and its overall impact."},
        "impact_level": {"type": "string", "enum": ["Low", "Medium", "High", "Critical"]},
        "affected_requirements": {
            "type": "array",
            "items": {
                "type": "object",
                "properties": {
                    "requirement": {"type": "string", "description": "The existing requirement/story affected, using its label from the source content."},
                    "impact_description": {"type": "string", "description": "Specifically how this change affects it (conflicts, extends, invalidates, etc.)."},
                },
                "required": ["requirement", "impact_description"],
            },
        },
        "recommended_actions": {
            "type": "array",
            "items": {"type": "string"},
            "description": "Concrete next steps, e.g. 'Re-confirm scope with sponsor', 'Update BRD section 4', 'Re-estimate story X'.",
        },
    },
    "required": ["summary", "impact_level", "affected_requirements", "recommended_actions"],
}


# --- Project State Management ---

def default_project():
    return {
        "description": "",
        "client": "",
        "status": "Planning",
        "documents": [],          # list of {"name","text","ext","added_at","char_count"}
        "extracted_text": "",     # last combined text analyzed in Elicitation tab
        "last_notes": "",         # last notes used in Elicitation tab
        "gap_analysis": None,
        "stories": [],
        "stories_drafted": 0,
        "documents_drafted": 0,
        "last_doc_draft": None,   # {"kind": "markdown"|"data_dictionary"|"asis_tobe", ...}
        "last_doc_type": None,
        "meeting_result": None,
        "workshop_prep": None,
        "test_cases": [],
        "test_cases_drafted": 0,
        "glossary": [],
        "glossary_terms_found": 0,
        "prioritization": None,
        "rtm_rows": [],
        "change_impact_history": [],   # list of {"request_text","result"} — most recent last
    }


def init_projects():
    if "projects" not in st.session_state:
        st.session_state["projects"] = {name: default_project() for name in STARTER_PROJECTS}
    if "current_project" not in st.session_state:
        st.session_state["current_project"] = STARTER_PROJECTS[0]


def get_current_project_name():
    return st.session_state.get("current_project")


def get_project():
    name = get_current_project_name()
    if name not in st.session_state["projects"]:
        st.session_state["projects"][name] = default_project()
    return st.session_state["projects"][name]


def add_doc_to_repo(proj, name, text, ext):
    existing_names = {d["name"] for d in proj["documents"]}
    if name in existing_names:
        proj["documents"] = [d for d in proj["documents"] if d["name"] != name]
    proj["documents"].append({
        "name": name,
        "text": text,
        "ext": ext,
        "added_at": datetime.now().strftime("%Y-%m-%d %H:%M"),
        "char_count": len(text),
    })


# --- API Helpers ---

def get_api_key():
    key = None
    try:
        key = st.secrets.get("ANTHROPIC_API_KEY")
    except Exception:
        key = None
    if not key:
        key = os.environ.get("ANTHROPIC_API_KEY")
    return key


@st.cache_resource
def get_client():
    api_key = get_api_key()
    if not api_key:
        st.error(
            "No Anthropic API key found. Add `ANTHROPIC_API_KEY` to this app's "
            "Secrets (on Streamlit Community Cloud: Settings → Secrets), or set it as an "
            "environment variable if running locally."
        )
        st.stop()
    # Deliberately no custom http_client here. An earlier version passed an
    # httpx.Client pinned to IPv4, to chase "Connection error" failures in the
    # nightly evals. That was the wrong diagnosis — the real cause was a
    # malformed ANTHROPIC_API_KEY secret — and the override then broke
    # production outright: newer anthropic SDKs validate this argument against
    # httpx2.Client and raise TypeError on an httpx.Client, so the type that
    # is correct depends on the SDK version installed. Letting the SDK build
    # its own client keeps this correct across versions.
    return anthropic.Anthropic(api_key=api_key)


def current_model():
    return st.session_state.get("model", DEFAULT_MODEL)


def call_text(system, user_prompt, max_tokens=1500, model=None):
    client = get_client()
    try:
        resp = client.messages.create(
            model=model or current_model(),
            max_tokens=max_tokens,
            system=system,
            messages=[{"role": "user", "content": user_prompt}],
        )
        return "".join(block.text for block in resp.content if block.type == "text")
    except Exception as e:
        st.error(f"AI request failed: {e}")
        return None


def call_chat(system, messages, max_tokens=800, model=None):
    client = get_client()
    try:
        resp = client.messages.create(
            model=model or current_model(),
            max_tokens=max_tokens,
            system=system,
            messages=messages,
        )
        return "".join(block.text for block in resp.content if block.type == "text")
    except Exception as e:
        st.error(f"AI request failed: {e}")
        return None


def call_structured(system, user_prompt, tool_name, tool_description, schema, max_tokens=2000, model=None):
    client = get_client()
    try:
        resp = client.messages.create(
            model=model or current_model(),
            max_tokens=max_tokens,
            system=system,
            messages=[{"role": "user", "content": user_prompt}],
            tools=[{"name": tool_name, "description": tool_description, "input_schema": schema}],
            tool_choice={"type": "tool", "name": tool_name},
        )
        for block in resp.content:
            if block.type == "tool_use" and block.name == tool_name:
                return block.input
        st.error("The AI response didn't include the expected structured data. Try again.")
        return None
    except Exception as e:
        st.error(f"AI request failed: {e}")
        return None


def call_structured_multimodal(system, content_blocks, tool_name, tool_description, schema, max_tokens=2000, model=None):
    """Like call_structured, but the user turn is a list of content blocks (images + text)
    instead of a plain string — used to send extracted images to Claude's vision model."""
    client = get_client()
    try:
        resp = client.messages.create(
            model=model or current_model(),
            max_tokens=max_tokens,
            system=system,
            messages=[{"role": "user", "content": content_blocks}],
            tools=[{"name": tool_name, "description": tool_description, "input_schema": schema}],
            tool_choice={"type": "tool", "name": tool_name},
        )
        for block in resp.content:
            if block.type == "tool_use" and block.name == tool_name:
                return block.input
        return None
    except Exception as e:
        st.warning(f"Image analysis skipped (AI request failed: {e}).")
        return None


# --- File Parsing (reading uploads) ---

MAX_CHARS = 15000
MAX_IMAGES_PER_DOC = 8       # cap vision calls per document — keeps cost/latency reasonable
IMAGE_MAX_DIMENSION = 1200   # px, longest side — plenty for screenshots/diagrams, keeps payload small


def _prepare_image_for_vision(image_bytes):
    """Decodes arbitrary image bytes, downsizes if needed, and re-encodes as JPEG base64
    for the vision API. Returns None if the bytes aren't a readable image (e.g. an EMF/WMF
    vector graphic that PIL can't open, or a corrupt embed) rather than raising."""
    try:
        img = Image.open(io.BytesIO(image_bytes))
        img = img.convert("RGB")
        if max(img.size) > IMAGE_MAX_DIMENSION:
            img.thumbnail((IMAGE_MAX_DIMENSION, IMAGE_MAX_DIMENSION))
        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=85)
        return base64.standard_b64encode(buf.getvalue()).decode("utf-8")
    except Exception:
        return None


def describe_images_with_vision(images, source_name):
    """Sends up to MAX_IMAGES_PER_DOC extracted images to Claude vision in a single call
    and returns a list of description strings (content_type: description), in the same
    order as the input images. Silently returns [] if there are no usable images or the
    vision call fails — image description is a bonus signal, not a hard requirement for
    the rest of the pipeline to work."""
    if not images:
        return []

    prepared = []
    for img_bytes in images[:MAX_IMAGES_PER_DOC]:
        b64 = _prepare_image_for_vision(img_bytes)
        if b64:
            prepared.append(b64)
    if not prepared:
        return []

    content_blocks = []
    for b64 in prepared:
        content_blocks.append({
            "type": "image",
            "source": {"type": "base64", "media_type": "image/jpeg", "data": b64},
        })
    content_blocks.append({
        "type": "text",
        "text": (
            f"These are {len(prepared)} image(s) embedded in the document '{source_name}', in order. "
            "Describe each one: if it's a screenshot of a UI/application, describe the screen, fields, "
            "buttons, and any visible text or labels — these often encode real requirements. If it's a "
            "diagram or flowchart, describe the steps/relationships shown. If it's a chart, describe the "
            "data and trend. Submit one entry per image, numbered in order."
        ),
    })

    result = call_structured_multimodal(
        "You are a business analyst assistant that reads screenshots and diagrams embedded in "
        "requirements documents and describes them precisely so their content can be analyzed as text.",
        content_blocks, "submit_image_analysis", "Submit descriptions of the provided images.",
        IMAGE_ANALYSIS_SCHEMA, max_tokens=2000,
    )
    if not result:
        return []
    descriptions = []
    for item in result.get("images", []):
        descriptions.append(f"[{item.get('content_type', 'Image')}] {item.get('description', '')}")
    return descriptions


def extract_docx_with_formatting(uploaded_file, describe_images=True):
    """Reads a .docx with formatting awareness instead of flattening to plain text.
    Bold runs are wrapped **like this**, struck-through runs ~~like this~~, and italic
    runs *like this* — markdown conventions the AI already understands, so emphasis and
    deprecated/removed content carry through as real signal. Inline comments (Word's
    actual comment feature, not Track Changes) are appended as a labeled section.
    Embedded images (screenshots, diagrams) are extracted and described via Claude vision.

    Known limitation: Word's Track Changes redline deletions/insertions (the dotted
    underline / strikethrough you see when "Show Markup" is on) are a different XML
    mechanism that python-docx doesn't expose at a usable level — only explicit
    strikethrough *formatting* applied to a run is captured here, not tracked changes.
    """
    doc = Document(uploaded_file)
    lines = []
    for para in doc.paragraphs:
        if not para.runs:
            if para.text.strip():
                lines.append(para.text)
            continue
        rendered = []
        for run in para.runs:
            t = run.text
            if not t:
                continue
            if run.font.strike:
                t = f"~~{t}~~"
            if run.bold:
                t = f"**{t}**"
            if run.italic:
                t = f"*{t}*"
            rendered.append(t)
        line = "".join(rendered)
        if line.strip():
            lines.append(line)
    body_text = "\n".join(lines)

    try:
        comments = list(doc.comments)
    except Exception:
        comments = []
    if comments:
        comment_lines = [f"- {c.author or 'Unknown reviewer'}: {c.text}" for c in comments if c.text and c.text.strip()]
        if comment_lines:
            body_text += "\n\n--- Reviewer Comments (from Word comments) ---\n" + "\n".join(comment_lines)

    if describe_images:
        image_bytes_list = []
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                try:
                    image_bytes_list.append(rel.target_part.blob)
                except Exception:
                    continue
        if image_bytes_list:
            with st.spinner(f"Reading {min(len(image_bytes_list), MAX_IMAGES_PER_DOC)} embedded image(s)..."):
                descriptions = describe_images_with_vision(image_bytes_list, uploaded_file.name)
            if descriptions:
                desc_lines = [f"{i+1}. {d}" for i, d in enumerate(descriptions)]
                body_text += "\n\n--- Images/Screenshots Found in Document (AI-described) ---\n" + "\n".join(desc_lines)

    return body_text


def extract_pdf_with_annotations(uploaded_file, describe_images=True):
    """Extracts PDF text plus any sticky-note/comment annotations, and describes any
    embedded images (scanned screenshots, diagrams) via Claude vision. Bold/strikethrough
    detection isn't attempted for PDFs — pypdf's text extraction doesn't expose per-
    character font styling, and strikethrough in a PDF is often just a drawn line
    rather than a text attribute, so it can't be reliably detected generically."""
    reader = pypdf.PdfReader(uploaded_file)
    if reader.is_encrypted:
        # Plenty of routine exports (DocuSign, SharePoint, Acrobat's "restrict
        # editing") are encrypted with an *empty* user password purely to limit
        # editing or printing — every PDF viewer opens these without prompting,
        # but pypdf refuses to read the pages until it's told to decrypt. Only
        # the empty password is tried: a PDF with a real password still fails,
        # as it should.
        try:
            unlocked = reader.decrypt("")
        except Exception:
            unlocked = False
        if not unlocked:
            raise ValueError(
                "This PDF is password-protected. Open it with the password, save an "
                "unprotected copy, and upload that instead."
            )
    pages = [page.extract_text() or "" for page in reader.pages]
    text = "\n".join(pages)
    if not text.strip():
        st.warning("No extractable text found in this PDF — it may be a scanned/image-only document.")

    annotation_lines = []
    for page in reader.pages:
        if not page.annotations:
            continue
        for a in page.annotations:
            obj = a.get_object()
            contents = obj.get("/Contents")
            if contents and str(contents).strip():
                annotation_lines.append(f"- {obj.get('/T', 'Unknown reviewer')}: {contents}")
    if annotation_lines:
        text += "\n\n--- Reviewer Comments (from PDF annotations) ---\n" + "\n".join(annotation_lines)

    if describe_images:
        image_bytes_list = []
        for page in reader.pages:
            try:
                for img in page.images:
                    image_bytes_list.append(img.data)
            except Exception:
                continue
        if image_bytes_list:
            with st.spinner(f"Reading {min(len(image_bytes_list), MAX_IMAGES_PER_DOC)} embedded image(s)..."):
                descriptions = describe_images_with_vision(image_bytes_list, uploaded_file.name)
            if descriptions:
                desc_lines = [f"{i+1}. {d}" for i, d in enumerate(descriptions)]
                text += "\n\n--- Images/Screenshots Found in Document (AI-described) ---\n" + "\n".join(desc_lines)

    return text


def _iter_pptx_shapes(shapes, _depth=0):
    """Yields a slide's shapes, descending into grouped ones.

    Grouping a screenshot with a caption box is a very common way to annotate a
    slide, and it nests both inside a GROUP shape — so a flat pass over
    slide.shapes skips them silently: no error, the picture and its text simply
    never reach the AI. The depth cap stops a pathologically nested file from
    recursing away; content below it is skipped rather than crashing."""
    if _depth > 10:
        return
    for shape in shapes:
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            yield from _iter_pptx_shapes(shape.shapes, _depth + 1)
        else:
            yield shape


def extract_pptx_with_formatting(uploaded_file, describe_images=True):
    """Reads a .pptx slide-by-slide: slide title/body text (with bold/italic markdown
    markers), speaker notes, and embedded images (screenshots often pasted into slides
    to illustrate a workflow) described via Claude vision."""
    prs = Presentation(uploaded_file)
    lines = []
    image_bytes_list = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_lines = [f"--- Slide {slide_num} ---"]
        for shape in _iter_pptx_shapes(slide.shapes):
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    image_bytes_list.append(shape.image.blob)
                except Exception:
                    pass
                continue
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                rendered = []
                for run in para.runs:
                    t = run.text
                    if not t:
                        continue
                    if run.font.bold:
                        t = f"**{t}**"
                    if run.font.italic:
                        t = f"*{t}*"
                    rendered.append(t)
                line = "".join(rendered)
                if line.strip():
                    slide_lines.append(line)
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text
            if notes_text.strip():
                slide_lines.append(f"[Speaker notes: {notes_text.strip()}]")
        if len(slide_lines) > 1:
            lines.extend(slide_lines)

    body_text = "\n".join(lines)

    if describe_images and image_bytes_list:
        with st.spinner(f"Reading {min(len(image_bytes_list), MAX_IMAGES_PER_DOC)} embedded image(s)..."):
            descriptions = describe_images_with_vision(image_bytes_list, uploaded_file.name)
        if descriptions:
            desc_lines = [f"{i+1}. {d}" for i, d in enumerate(descriptions)]
            body_text += "\n\n--- Images/Screenshots Found in Presentation (AI-described) ---\n" + "\n".join(desc_lines)

    return body_text


MAX_UPLOAD_BYTES = 20 * 1024 * 1024  # cap uploaded files before they're parsed/sent to vision
MAX_UNCOMPRESSED_BYTES = 200 * 1024 * 1024  # cap what a zip-backed Office file expands to
_ZIP_BACKED_EXTS = ("docx", "pptx", "xlsx")


def _assert_safe_archive(uploaded_file, ext):
    """Guards against a decompression bomb in a zip-backed Office file.

    .docx/.pptx/.xlsx are zip archives, so MAX_UPLOAD_BYTES only caps the
    *compressed* bytes on the way in. python-docx / python-pptx / openpyxl all
    unzip internally, so a small archive can declare gigabytes of content and
    exhaust memory before any of the app's other limits apply. Reading the
    archive's central directory gives the declared sizes without decompressing
    anything, so this costs almost nothing.

    This trusts the sizes the archive declares. A hand-crafted zip can
    understate them, so treat it as a guard against the realistic case rather
    than a hard guarantee.
    """
    uploaded_file.seek(0)
    try:
        with zipfile.ZipFile(uploaded_file) as zf:
            total = sum(info.file_size for info in zf.infolist())
    except zipfile.BadZipFile:
        raise ValueError(
            f"This .{ext} file isn't a readable Office document — it may be corrupted."
        ) from None
    finally:
        uploaded_file.seek(0)

    if total > MAX_UNCOMPRESSED_BYTES:
        raise ValueError(
            f"This .{ext} file expands to about {total // (1024*1024)}MB when opened, "
            f"over the {MAX_UNCOMPRESSED_BYTES // (1024*1024)}MB limit — it can't be "
            "processed safely."
        )


def extract_text_from_upload(uploaded_file):
    name = uploaded_file.name
    ext = name.split(".")[-1].lower()
    if uploaded_file.size > MAX_UPLOAD_BYTES:
        st.error(f"This file is larger than {MAX_UPLOAD_BYTES // (1024*1024)}MB — too large to process.")
        return ""
    if ext in _ZIP_BACKED_EXTS:
        try:
            _assert_safe_archive(uploaded_file, ext)
        except ValueError as e:
            st.error(str(e))
            return ""
    uploaded_file.seek(0)
    try:
        if ext == "txt":
            raw = uploaded_file.read()
            if raw[:2] in (b"\xff\xfe", b"\xfe\xff"):
                text = raw.decode("utf-16")
            else:
                try:
                    text = raw.decode("utf-8-sig")
                except UnicodeDecodeError:
                    text = raw.decode("cp1252", errors="ignore")
        elif ext == "pdf":
            text = extract_pdf_with_annotations(uploaded_file)
        elif ext == "docx":
            text = extract_docx_with_formatting(uploaded_file)
        elif ext == "pptx":
            text = extract_pptx_with_formatting(uploaded_file)
        elif ext == "csv":
            df = pd.read_csv(uploaded_file)
            text = df.to_string(index=False) if not df.empty else ""
        elif ext == "xlsx":
            xls = pd.ExcelFile(uploaded_file)
            parts = []
            for sheet in xls.sheet_names:
                df = xls.parse(sheet)
                if not df.empty:
                    parts.append(f"--- Sheet: {sheet} ---\n{df.to_string(index=False)}")
            text = "\n\n".join(parts)
        else:
            st.error(f"Unsupported file type: .{ext}")
            return ""
    except ValueError as e:
        st.error(str(e))
        return ""
    except Exception as e:
        st.error(f"Couldn't read this file: {e}")
        return ""

    if ext != "pdf" and not text.strip():
        if ext in ("docx", "pptx"):
            st.warning(f"No extractable text found in this .{ext} file — it may be empty or image-only.")
        else:
            st.warning(f"No extractable text found in this .{ext} file — it may be empty.")
    return text


MAX_FETCH_URL_BYTES = 5 * 1024 * 1024  # cap page size before it's pulled into memory/parsed
MAX_FETCH_REDIRECTS = 5


def _assert_public_url(url):
    """Rejects a URL that doesn't point at a public internet host.

    Source URLs are user-supplied but fetched *server-side*, so without this check
    the fetcher could be aimed at addresses only the server can reach — cloud
    instance-metadata endpoints, admin panels on the host network, localhost
    services — and their responses read back through the app. Every address the
    hostname resolves to is checked, so a public-looking hostname with a private
    A record is rejected too.

    Limitation: this validates at resolve time. A hostname whose DNS answer
    changes between this check and the actual connection (DNS rebinding) is not
    defeated by this alone; that needs connection-level IP pinning, which is a
    bigger change than belongs here.
    """
    parsed = urlparse(url)
    if parsed.scheme not in ("http", "https"):
        raise ValueError(
            f"Only http:// and https:// URLs can be fetched (got "
            f"'{parsed.scheme or 'no scheme'}')."
        )
    host = parsed.hostname
    if not host:
        raise ValueError("That doesn't look like a valid URL — no hostname found.")

    port = parsed.port or (443 if parsed.scheme == "https" else 80)
    try:
        infos = socket.getaddrinfo(host, port, proto=socket.IPPROTO_TCP)
    except socket.gaierror:
        raise ValueError(f"Couldn't resolve '{host}' — check the address and try again.") from None

    for info in infos:
        ip = ipaddress.ip_address(info[4][0])
        if ip.version == 6 and ip.ipv4_mapped:
            ip = ip.ipv4_mapped
        if (ip.is_private or ip.is_loopback or ip.is_link_local
                or ip.is_reserved or ip.is_multicast or ip.is_unspecified):
            raise ValueError(
                f"'{host}' resolves to a private or internal address ({ip}). "
                "Only public webpages can be fetched."
            )


def fetch_url_text(url, timeout=10):
    """Fetches a public webpage and returns its visible text plus a best-guess title.
    Only handles publicly-accessible pages — see the 'Pages that require a login' note
    in the Project & Documents tab for why authenticated fetching isn't implemented."""
    headers = {"User-Agent": "Mozilla/5.0 (compatible; ScopeForge-DocBot/1.0)"}

    # Redirects are followed by hand so every hop gets the same public-address
    # check as the original URL — otherwise a public URL could bounce the fetch
    # to an internal one and sidestep the check entirely.
    current = url
    for _ in range(MAX_FETCH_REDIRECTS + 1):
        _assert_public_url(current)
        resp = requests.get(current, headers=headers, timeout=timeout,
                            stream=True, allow_redirects=False)
        if 300 <= resp.status_code < 400:
            location = resp.headers.get("Location")
            resp.close()
            if not location:
                raise ValueError("This URL redirected without saying where to.")
            current = urljoin(current, location)
            continue
        break
    else:
        raise ValueError(
            f"This URL redirected more than {MAX_FETCH_REDIRECTS} times — giving up."
        )

    resp.raise_for_status()

    chunks = []
    total = 0
    for chunk in resp.iter_content(chunk_size=65536):
        total += len(chunk)
        if total > MAX_FETCH_URL_BYTES:
            resp.close()
            raise ValueError(f"This page is larger than {MAX_FETCH_URL_BYTES // (1024*1024)}MB — too large to fetch.")
        chunks.append(chunk)
    resp.encoding = resp.encoding or "utf-8"
    try:
        body = b"".join(chunks).decode(resp.encoding, errors="ignore")
    except (LookupError, TypeError):
        body = b"".join(chunks).decode("utf-8", errors="ignore")

    content_type = resp.headers.get("Content-Type", "").lower()
    if "text/plain" in content_type:
        return body, url

    soup = BeautifulSoup(body, "html.parser")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    title = soup.title.string.strip() if (soup.title and soup.title.string) else url
    raw = soup.get_text(separator="\n")
    cleaned = "\n".join(line.strip() for line in raw.splitlines() if line.strip())
    return cleaned, title


def truncate(text, limit=MAX_CHARS):
    if len(text) > limit:
        return text[:limit], True
    return text, False


def _blank_or_value(v):
    """Guards against pd.isna() raising on a non-scalar cell (e.g. a list/dict that
    slipped through from an AI response field expected to be a plain string) — such a
    value is never actually NaN, so it's stringified instead of being checked, which
    also keeps it writable as a single Excel/Word cell value."""
    if isinstance(v, (list, dict)):
        return str(v)
    return "" if pd.isna(v) else v


# --- File Building (generating downloads) ---

def build_xlsx_from_df(sheet_name, df):
    wb = Workbook()
    ws = wb.active
    ws.title = (sheet_name or "Sheet1")[:31]
    ws.append([str(c) for c in df.columns])
    for cell in ws[1]:
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill(start_color="1F2333", end_color="1F2333", fill_type="solid")
    for _, row in df.iterrows():
        ws.append([_blank_or_value(v) for v in row.tolist()])
    for i in range(1, len(df.columns) + 1):
        ws.column_dimensions[ws.cell(row=1, column=i).column_letter].width = 28
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def build_docx_table_from_df(title, df):
    doc = Document()
    doc.add_heading(title, level=0)
    table = doc.add_table(rows=1, cols=len(df.columns))
    table.style = "Light Grid Accent 1"
    hdr_cells = table.rows[0].cells
    for i, col in enumerate(df.columns):
        hdr_cells[i].text = str(col)
    for _, row in df.iterrows():
        cells = table.add_row().cells
        for i, v in enumerate(row.tolist()):
            cells[i].text = str(_blank_or_value(v))
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def build_docx_from_markdown(title, markdown_text):
    doc = Document()
    doc.add_heading(title, level=0)
    for raw_line in markdown_text.split("\n"):
        line = raw_line.strip()
        if not line:
            continue
        if line.startswith("#### "):
            doc.add_heading(line[5:], level=4)
        elif line.startswith("### "):
            doc.add_heading(line[4:], level=3)
        elif line.startswith("## "):
            doc.add_heading(line[3:], level=2)
        elif line.startswith("# "):
            doc.add_heading(line[2:], level=1)
        elif re.match(r"^[-*]\s+", line):
            doc.add_paragraph(re.sub(r"^[-*]\s+", "", line), style="List Bullet")
        elif re.match(r"^\d+\.\s+", line):
            doc.add_paragraph(re.sub(r"^\d+\.\s+", "", line), style="List Number")
        else:
            doc.add_paragraph(line.replace("**", "").replace("*", "").replace("~~", ""))
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


VISIO_COLUMNS = ["Process Step ID", "Process Step Description", "Next Step ID", "Connector Label", "Shape Type"]


def build_visio_dataviz_xlsx(rows):
    """Builds an Excel file in Microsoft's documented Data Visualizer 'Basic Flowchart'
    column format. Import this into Visio (Data > Create from Data / Data Visualizer
    template) to auto-generate an editable, real Visio diagram — this is the reliable,
    well-documented path; authoring a raw .vsdx binary from scratch isn't something that
    can be done robustly without Visio itself or a paid SDK."""
    wb = Workbook()
    ws = wb.active
    ws.title = "Process Map"
    ws.append(VISIO_COLUMNS)
    for cell in ws[1]:
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill(start_color="1F2333", end_color="1F2333", fill_type="solid")
    for r in rows:
        ws.append([
            str(r.get("step_id", "")),
            r.get("process_step", ""),
            str(r.get("next_step_id", "")),
            "",
            r.get("shape_type", "Process"),
        ])
    # Force the ID columns to explicit text format — otherwise Excel may auto-interpret
    # values like "010" as the number 10, which would break step matching on import.
    for row in ws.iter_rows(min_row=2, min_col=1, max_col=1):
        row[0].number_format = "@"
    for row in ws.iter_rows(min_row=2, min_col=3, max_col=3):
        row[0].number_format = "@"
    for i in range(1, len(VISIO_COLUMNS) + 1):
        ws.column_dimensions[ws.cell(row=1, column=i).column_letter].width = 28
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# --- AI Task Functions ---

FORMATTING_GUIDANCE = (
    "\n\nSource content carried over from Word documents may include markdown-style "
    "formatting that reflects the original document's markup: **bold** text indicates "
    "emphasis or a critical/non-negotiable requirement; ~~struck-through~~ text indicates "
    "content that was removed, deprecated, or rejected — treat it as historical/As-Is "
    "context, not a current requirement, unless asked specifically about prior versions. "
    "A 'Reviewer Comments' section (from Word comments or PDF annotations) contains "
    "stakeholder/reviewer feedback — factor it into your analysis as you would any other "
    "stated constraint or concern."
)


def analyze_gaps(text, notes=None):
    truncated_text, was_truncated = truncate(text)
    if was_truncated:
        st.caption(f"Document content was long — analyzing the first {MAX_CHARS:,} characters.")
    system = (
        "You are a senior business analyst performing a requirements quality review. "
        "Carefully read the provided notes/documents and identify ambiguities, missing "
        "non-functional requirements (performance, security, availability, etc.), stakeholder "
        "conflicts, and scope risks. Be specific and ground every finding in something actually "
        "present (or notably absent) in the text — do not invent details. If the text is sparse, "
        "it's fine to return fewer findings and a lower risk score."
    ) + FORMATTING_GUIDANCE
    user_prompt = f"Analyze the following content:\n\n{truncated_text}"
    if notes and notes.strip():
        user_prompt += (
            "\n\n---\nThe business analyst has also provided the following additional notes/context. "
            "You MUST factor these into your analysis — they may resolve an apparent gap, introduce "
            "a new constraint or decision, or point you toward something specific to scrutinize:\n"
            f"{notes.strip()}"
        )
    return call_structured(
        system, user_prompt, "submit_gap_analysis",
        "Submit the structured requirements gap analysis.", GAP_ANALYSIS_SCHEMA,
    )


def generate_document(doc_type, context_text, user_suggestion):
    context_text, was_truncated = truncate(context_text)
    if was_truncated:
        st.caption(f"Source content was long — using the first {MAX_CHARS:,} characters.")
    system = (
        "You are a senior business analyst drafting professional project documentation. "
        "Write a well-structured, realistic first draft in Markdown. Use clear section headers. "
        "Where the source content doesn't cover something, write '[Needs stakeholder input]' "
        "rather than inventing specifics."
    ) + FORMATTING_GUIDANCE
    user_prompt = (
        f"Document type to draft: {doc_type}\n\n"
        f"Special instructions / focus areas from the analyst:\n{user_suggestion}\n\n"
        f"Source content / requirements notes to base this on:\n"
        f"{context_text if context_text.strip() else '(No source content provided — draft a generic template with clear placeholder sections.)'}"
    )
    return call_text(system, user_prompt, max_tokens=2500)


def generate_data_dictionary(context_text, user_suggestion):
    context_text, was_truncated = truncate(context_text)
    if was_truncated:
        st.caption(f"Source content was long — using the first {MAX_CHARS:,} characters.")
    system = (
        "You are a senior business analyst building a data dictionary. Identify every distinct "
        "data field/entity attribute implied by the source content and document it. Only include "
        "fields actually supported by the source content."
    ) + FORMATTING_GUIDANCE
    user_prompt = (
        f"Special instructions / focus areas:\n{user_suggestion}\n\n"
        f"Source content:\n{context_text if context_text.strip() else '(No source content provided.)'}"
    )
    result = call_structured(
        system, user_prompt, "submit_data_dictionary",
        "Submit the structured data dictionary.", DATA_DICT_SCHEMA, max_tokens=3000,
    )
    return result.get("fields", []) if result else []


def generate_asis_tobe(context_text, user_suggestion):
    context_text, was_truncated = truncate(context_text)
    if was_truncated:
        st.caption(f"Source content was long — using the first {MAX_CHARS:,} characters.")
    system = (
        "You are a senior business analyst mapping a business process for an As-Is / To-Be "
        "analysis. Break the process into sequential steps. For each step, describe the current "
        "(As-Is) state, the proposed future (To-Be) state, and the specific gap/change needed to "
        "get from one to the other. Assign a shape_type for flowchart purposes (Start/Process/"
        "Decision/End) and a next_step_id describing what step(s) follow (comma-separated, no "
        "spaces, blank for the End step). Only build steps actually supported by the source content. "
        "Pay close attention to struck-through text in the source — it usually marks a requirement "
        "or process step that was cut, so it belongs in the As-Is description (as something that "
        "existed/was proposed) but generally should NOT carry into the To-Be description unless the "
        "source clearly indicates it's being reinstated."
    ) + FORMATTING_GUIDANCE
    user_prompt = (
        f"Special instructions / focus areas:\n{user_suggestion}\n\n"
        f"Source content describing the process:\n{context_text if context_text.strip() else '(No source content provided.)'}"
    )
    result = call_structured(
        system, user_prompt, "submit_asis_tobe",
        "Submit the structured As-Is / To-Be process map.", ASIS_TOBE_SCHEMA, max_tokens=3000,
    )
    return result.get("steps", []) if result else []


def generate_stories(source_text):
    source_text, was_truncated = truncate(source_text)
    if was_truncated:
        st.caption(f"Source content was long — using the first {MAX_CHARS:,} characters.")
    system = (
        "You are a senior business analyst converting requirements into an Agile backlog. "
        "For each distinct requirement or need in the source text, write one user story in the "
        "format 'As a <role>, I want <capability>, so that <benefit>', plus Gherkin-style "
        "acceptance criteria (GIVEN/WHEN/THEN). Only create stories that are actually supported "
        "by the source text."
    ) + FORMATTING_GUIDANCE
    user_prompt = f"Source requirements/notes:\n\n{source_text}"
    result = call_structured(
        system, user_prompt, "submit_stories",
        "Submit the generated user stories and acceptance criteria.", STORY_SCHEMA, max_tokens=2500,
    )
    return result.get("stories", []) if result else []


def generate_test_cases(stories):
    """Builds test cases from already-generated user stories/acceptance criteria."""
    story_lines = []
    for s in stories:
        story_lines.append(
            f"Requirement: {s.get('requirement') or ''}\n"
            f"User Story: {s.get('user_story') or ''}\n"
            f"Acceptance Criteria: {s.get('acceptance_criteria') or ''}"
        )
    source_text, was_truncated = truncate("\n\n".join(story_lines))
    if was_truncated:
        st.caption(f"Story list was long — using the first {MAX_CHARS:,} characters.")
    system = (
        "You are a senior QA/business analyst writing test cases from user stories and their "
        "Gherkin acceptance criteria. For each story, write one or more test cases that verify "
        "its acceptance criteria — cover the happy path and, where the acceptance criteria imply "
        "one, at least one edge case. Only test what the acceptance criteria actually specify; "
        "don't invent behavior."
    )
    user_prompt = f"User stories and acceptance criteria to derive test cases from:\n\n{source_text}"
    result = call_structured(
        system, user_prompt, "submit_test_cases",
        "Submit the generated test cases.", TEST_CASE_SCHEMA, max_tokens=3000,
    )
    return result.get("test_cases", []) if result else []


def generate_glossary(context_text):
    context_text, was_truncated = truncate(context_text)
    if was_truncated:
        st.caption(f"Source content was long — using the first {MAX_CHARS:,} characters.")
    system = (
        "You are a senior business analyst building a business glossary. Read the source content "
        "and extract domain-specific terms, acronyms, and system/product names that a new team "
        "member would need explained — skip common English words. Ground every definition in how "
        "the term is actually used in the source; don't invent definitions for terms not present."
    ) + FORMATTING_GUIDANCE
    user_prompt = f"Source content:\n\n{context_text if context_text.strip() else '(No source content provided.)'}"
    result = call_structured(
        system, user_prompt, "submit_glossary",
        "Submit the extracted business glossary.", GLOSSARY_SCHEMA, max_tokens=2500,
    )
    return result.get("terms", []) if result else []


def generate_prioritization(context_text):
    context_text, was_truncated = truncate(context_text)
    if was_truncated:
        st.caption(f"Source content was long — using the first {MAX_CHARS:,} characters.")
    system = (
        "You are a senior business analyst facilitating MoSCoW prioritization. Read the source "
        "content, identify each distinct requirement or need, and classify it as Must Have, Should "
        "Have, Could Have, or Won't Have (this time), with a one-sentence rationale grounded in the "
        "source content (e.g., stated urgency, dependency, regulatory need, or explicit stakeholder "
        "priority). If priority isn't stated explicitly, use reasonable business judgment and say so "
        "in the rationale."
    ) + FORMATTING_GUIDANCE
    user_prompt = f"Source content:\n\n{context_text if context_text.strip() else '(No source content provided.)'}"
    result = call_structured(
        system, user_prompt, "submit_prioritization",
        "Submit the MoSCoW prioritization.", PRIORITIZATION_SCHEMA, max_tokens=2500,
    )
    return result.get("items", []) if result else []


def generate_workshop_prep(project_description, focus_area, existing_context):
    existing_context, was_truncated = truncate(existing_context or "")
    if was_truncated:
        st.caption(f"Existing project context was long — using the first {MAX_CHARS:,} characters.")
    system = (
        "You are a senior business analyst preparing for a stakeholder elicitation workshop. "
        "Produce a tight, time-boxed agenda and a targeted list of open-ended elicitation "
        "questions grouped by category (e.g., Scope, Process, Data, Constraints, Success Criteria). "
        "Base the questions on what's already known from the project description/context — probe "
        "specifically for what's missing or ambiguous, don't ask generic questions that the context "
        "already answers."
    )
    user_prompt = (
        f"Project description:\n{project_description or '(none provided)'}\n\n"
        f"Workshop focus area:\n{focus_area or '(general requirements elicitation)'}\n\n"
        f"Existing project context/documents (if any):\n{existing_context or '(none)'}"
    )
    return call_structured(
        system, user_prompt, "submit_workshop_prep",
        "Submit the workshop agenda and question list.", WORKSHOP_PREP_SCHEMA, max_tokens=2000,
    )


def _label_overlap(a, b):
    """Fuzzy match score (0-1) between two short AI-generated labels, by word overlap.
    Three separate AI calls (stories, test cases, prioritization) each phrase the same
    underlying requirement in their own words — exact/substring string matching misses
    most real links, so this compares word sets instead of characters."""
    a_words = set(re.findall(r"[a-z0-9]+", a.lower()))
    b_words = set(re.findall(r"[a-z0-9]+", b.lower()))
    if not a_words or not b_words:
        return 0.0
    return len(a_words & b_words) / min(len(a_words), len(b_words))


RTM_MATCH_THRESHOLD = 0.4


def build_rtm_rows(proj):
    """Auto-builds Requirements Traceability Matrix rows from what's already been generated
    elsewhere in the project (stories, test cases, MoSCoW prioritization) — no separate AI
    call needed since the data already exists; this just links it together. Matching is
    fuzzy (word overlap, see _label_overlap) rather than exact, since the three source
    lists come from independent AI calls that won't phrase a requirement identically. It's
    still not perfect — that's why the RTM table stays editable, so the BA can fix/add
    links by hand."""
    stories = proj.get("stories", [])
    test_cases = proj.get("test_cases", [])
    prioritization = proj.get("prioritization") or []

    rows = []
    for s in stories:
        req_label = s.get("requirement", "")
        related_tests = [
            tc.get("test_id", "") for tc in test_cases
            if req_label and _label_overlap(req_label, tc.get("related_story", "")) >= RTM_MATCH_THRESHOLD
        ]
        prio = max(
            prioritization,
            key=lambda p: _label_overlap(req_label, p.get("requirement", "")),
            default=None,
        )
        if prio and _label_overlap(req_label, prio.get("requirement", "")) < RTM_MATCH_THRESHOLD:
            prio = None
        rows.append({
            "requirement": req_label,
            "user_story": s.get("user_story", ""),
            "test_cases": ", ".join(t for t in related_tests if t),
            "priority": prio.get("moscow_category", "") if prio else "",
            "status": "Drafted",
        })
    return rows


def generate_change_impact(change_request_text, existing_context):
    existing_context, was_truncated = truncate(existing_context)
    if was_truncated:
        st.caption(f"Existing project content was long — using the first {MAX_CHARS:,} characters.")
    system = (
        "You are a senior business analyst performing a change impact assessment. A change "
        "request has come in against a project with existing requirements/stories. Identify "
        "which existing requirements or stories it affects and how (conflicts with, extends, "
        "invalidates, adds a dependency to, etc.), rate the overall impact level, and recommend "
        "concrete next actions. Ground every affected item in something actually present in the "
        "existing content — don't invent requirements that aren't there. If nothing in the "
        "existing content is affected, say so and return an empty affected_requirements list."
    )
    user_prompt = (
        f"Change request:\n{change_request_text}\n\n"
        f"Existing project requirements/stories/documentation:\n"
        f"{existing_context if existing_context.strip() else '(No existing project content available — assess the change request on its own.)'}"
    )
    return call_structured(
        system, user_prompt, "submit_change_impact",
        "Submit the change impact assessment.", CHANGE_IMPACT_SCHEMA, max_tokens=2000,
    )


def process_meeting(transcript_text):
    truncated_text, was_truncated = truncate(transcript_text)
    if was_truncated:
        st.caption(f"Transcript was long — analyzing the first {MAX_CHARS:,} characters.")
    system = (
        "You are an assistant that turns raw meeting transcripts into structured minutes. "
        "Extract a concise executive summary, key decisions, and action items with an owner "
        "and due date if stated. Use 'Unassigned' / 'Not specified' rather than guessing."
    ) + FORMATTING_GUIDANCE
    user_prompt = f"Meeting transcript:\n\n{truncated_text}"
    return call_structured(
        system, user_prompt, "submit_meeting_minutes",
        "Submit the structured meeting minutes.", MEETING_SCHEMA,
    )


MAX_CHAT_MESSAGES_SENT = 20  # roughly ten back-and-forth exchanges of context
MAX_CHAT_MESSAGE_CHARS = 6000  # one pasted wall of text shouldn't crowd out the rest


def _recent_chat_messages(history):
    """Bounds the conversation handed to the API, without touching what's on screen.

    chat_history lives for the whole session and every turn resends all of it, so
    a long working session steadily increases cost and eventually runs into the
    model's context limit. Trimming here rather than in session_state keeps the
    user's full transcript visible while capping what actually goes over the wire.
    """
    trimmed = []
    for msg in history[-MAX_CHAT_MESSAGES_SENT:]:
        content = msg["content"]
        if len(content) > MAX_CHAT_MESSAGE_CHARS:
            content = content[:MAX_CHAT_MESSAGE_CHARS] + "\n\n[...truncated for length...]"
        trimmed.append({"role": msg["role"], "content": content})
    # The API needs the first message to come from the user, and a window taken
    # out of the middle of a conversation can easily start on an assistant reply.
    while trimmed and trimmed[0]["role"] != "user":
        trimmed.pop(0)
    return trimmed


def chat_with_bot(history):
    proj = get_project()
    proj_name = get_current_project_name()
    context_note = f"\n\nThe user is currently working in the '{proj_name}' project."
    if proj.get("description"):
        context_note += f" Project description: {proj['description'][:500]}"
    if proj.get("documents"):
        doc_names = ", ".join(d["name"] for d in proj["documents"][:10])
        context_note += f" Documents in this project's repository: {doc_names}."
    return call_chat(CHATBOT_SYSTEM_PROMPT + context_note,
                     _recent_chat_messages(history), max_tokens=800)


# --- Role-Specific Functions ---

def render_dashboard(proj, cp):
    section_header("BA Dashboard", f"Summary view for {cp} in this session.")

    col1, col2, col3, col4 = st.columns(4)
    col1.metric("User Stories Drafted", proj.get("stories_drafted", 0))
    col2.metric("Documents Drafted", proj.get("documents_drafted", 0))
    col3.metric("Repository Documents", len(proj.get("documents", [])))

    gap_result = proj.get("gap_analysis")
    open_gaps = len(gap_result.get("open_questions", [])) if gap_result else 0
    col4.metric("Open Gaps (latest analysis)", open_gaps)

    col5, col6, col7, col8 = st.columns(4)
    col5.metric("Test Cases Drafted", proj.get("test_cases_drafted", 0))
    col6.metric("Glossary Terms", proj.get("glossary_terms_found", 0))
    prioritized = proj.get("prioritization") or []
    must_haves = len([i for i in prioritized if i.get("moscow_category") == "Must Have"])
    col7.metric("Must-Have Requirements", must_haves)
    col8.metric("Workshop Preps Drafted", 1 if proj.get("workshop_prep") else 0)

    col9, col10 = st.columns(2)
    col9.metric("RTM Rows Tracked", len(proj.get("rtm_rows", [])))
    col10.metric("Change Requests Analyzed", len(proj.get("change_impact_history", [])))

    st.caption(
        "These metrics reflect activity for this project in your current browser session only. "
        "Tracking activity persistently across sessions/users would require a database backend."
    )


def ba_module():
    st.markdown(
        "<div style='font-size:1.7rem; font-weight:700;'>Strategic Requirements Hub</div>",
        unsafe_allow_html=True,
    )
    st.caption("AI-augmented tools for elicitation, documentation, and Agile workflow.")

    cp = get_current_project_name()
    proj = get_project()

    view = st.radio(
        "View", ["Workspace", "Dashboard"],
        horizontal=True, key=f"ba_view_{cp}", label_visibility="collapsed",
    )
    st.markdown("<hr style='margin-top:0.2rem;margin-bottom:1rem;'>", unsafe_allow_html=True)

    if view == "Dashboard":
        render_dashboard(proj, cp)
        return

    tab0, tab1, tab2, tab3, tab4, tab5 = st.tabs([
        "Project & Documents",
        "Meeting Intelligence & Actionizer",
        "Elicitation Analysis & Gap Detector",
        "Documentation Generator",
        "Agile Story & Backlog Creator",
        "Traceability & Change Impact",
    ])

    # --- Tab 0: Project & Documents ---
    with tab0:
        section_header("Project Information & Document Repository", f"Active project: {cp}")

        col1, col2 = st.columns(2)
        with col1:
            proj["client"] = st.text_input("Client / Stakeholder", value=proj.get("client", ""), key=f"client_{cp}")
        with col2:
            current_status = proj.get("status", "Planning")
            proj["status"] = st.selectbox(
                "Status", PROJECT_STATUSES,
                index=PROJECT_STATUSES.index(current_status) if current_status in PROJECT_STATUSES else 0,
                key=f"status_{cp}",
            )

        proj["description"] = st.text_area(
            "Project Description / Background", value=proj.get("description", ""), height=120, key=f"desc_{cp}",
        )

        st.markdown("---")
        st.markdown("#### Create a New Project")
        with st.form(key="new_project_form", clear_on_submit=True):
            new_name = st.text_input("New Project Name")
            new_desc = st.text_area("Description (optional)", height=80)
            submitted = st.form_submit_button("Create Project")
            if submitted:
                if not new_name.strip():
                    st.warning("Give the project a name.")
                elif new_name in st.session_state["projects"]:
                    st.warning("A project with this name already exists.")
                else:
                    new_proj = default_project()
                    new_proj["description"] = new_desc
                    st.session_state["projects"][new_name] = new_proj
                    st.session_state["pending_project_switch"] = new_name
                    st.success(f"Created project '{new_name}'.")
                    st.rerun()

        st.markdown("---")
        st.markdown("#### Document Repository")
        st.caption(
            "Documents added here are available as context to the Elicitation Analysis, "
            "Documentation Generator, and Story Creator tabs for this project. Word documents "
            "(.docx) are read with formatting awareness: **bold** text, ~~struck-through~~ text, "
            "and inline comments are preserved as signals for the AI — not flattened away. PDFs "
            "carry over sticky-note comments too. (Word's Track Changes redlines are a separate "
            "mechanism this can't see — only explicit strikethrough formatting and regular "
            "comments are captured.)"
        )

        repo_files = st.file_uploader(
            "Add documents to this project's repository",
            type=['txt', 'pdf', 'docx', 'pptx', 'xlsx', 'csv'],
            accept_multiple_files=True,
            key=f"repo_uploader_{cp}",
        )
        if st.button("Add to Repository", key=f"add_repo_btn_{cp}"):
            if not repo_files:
                st.warning("Choose at least one file first.")
            else:
                added = 0
                for f in repo_files:
                    text = extract_text_from_upload(f)
                    if text.strip():
                        add_doc_to_repo(proj, f.name, text, f.name.split(".")[-1].lower())
                        added += 1
                if added:
                    st.success(f"Added {added} document(s) to the repository.")
                    st.rerun()
                else:
                    st.warning("No documents were added — none of the selected file(s) had readable text.")

        if proj["documents"]:
            st.markdown(f"**{len(proj['documents'])} document(s) in repository:**")
            for i, doc in enumerate(proj["documents"]):
                c1, c2, c3, c4 = st.columns([4, 1, 2, 1])
                c1.write(doc["name"])
                c2.write(doc["ext"].upper())
                c3.write(f"{doc['char_count']:,} chars · added {doc['added_at']}")
                if c4.button("Remove", key=f"del_doc_{cp}_{i}"):
                    proj["documents"].pop(i)
                    st.rerun()
        else:
            st.info("No documents yet. Upload files above to build this project's repository.")

        st.markdown("---")
        st.markdown("#### Source URLs")
        st.caption(
            "Pull a public webpage's text straight into this project's repository, same as an "
            "uploaded file."
        )
        url_col1, url_col2 = st.columns([4, 1])
        with url_col1:
            url_input = st.text_input(
                "Public webpage URL", key=f"url_input_{cp}",
                placeholder="https://example.com/project-charter", label_visibility="collapsed",
            )
        with url_col2:
            fetch_clicked = st.button("Fetch URL", key=f"fetch_url_btn_{cp}", use_container_width=True)

        if fetch_clicked:
            if not url_input.strip():
                st.warning("Enter a URL first.")
            else:
                with st.spinner(f"Fetching {url_input.strip()}..."):
                    try:
                        text, title = fetch_url_text(url_input.strip())
                        if text.strip():
                            doc_name = f"[URL] {title or url_input.strip()}"
                            add_doc_to_repo(proj, doc_name, text, "url")
                            st.success(f"Added '{doc_name}' to the repository ({len(text):,} characters).")
                            st.rerun()
                        else:
                            st.warning("Fetched the page, but found no readable text on it.")
                    except requests.exceptions.HTTPError as e:
                        code = e.response.status_code if e.response is not None else None
                        if code in (401, 403):
                            st.error(
                                f"This page returned a {code} error — it likely requires you to be "
                                "logged in. See the note below."
                            )
                        else:
                            st.error(f"Couldn't fetch this URL ({code or 'HTTP error'}).")
                    except requests.exceptions.RequestException as e:
                        st.error(f"Couldn't reach this URL: {e}")
                    except ValueError as e:
                        st.error(str(e))

        with st.expander("Pages that require a login"):
            st.markdown(
                "This intentionally does **not** ask for a username/password to log into a site "
                "automatically. Two reasons:\n\n"
                "1. **Security** — handling another site's credentials inside this app means they "
                "pass through its server-side code. That's a real exposure surface for something "
                "that's supposed to be a lightweight internal tool, even if nothing is stored.\n"
                "2. **Reliability** — login flows vary enormously (SSO, MFA, CAPTCHAs, JavaScript-"
                "rendered forms). A generic auto-login would break on most real sites anyway, "
                "which isn't worth the security tradeoff above.\n\n"
                "The reliable workaround: open the page in your own browser while logged in, then "
                "either save it as a PDF/Word file or copy the text, and upload it using the "
                "Document Repository above. Same end result, no credentials ever touch this app."
            )

        st.markdown("---")
        with st.expander("Business Glossary", expanded=False):
            st.caption(
                "Scans this project's repository documents and pulls out domain terms, acronyms, "
                "and system names with plain-language definitions — handy for onboarding new "
                "team members or aligning stakeholders on vocabulary."
            )
            if not proj["documents"]:
                st.info("Add documents to the repository above first.")
            else:
                all_repo_doc_names = [d["name"] for d in proj["documents"]]
                glossary_doc_names = st.multiselect(
                    "Build glossary from:", all_repo_doc_names, default=all_repo_doc_names, key=f"glossary_repo_select_{cp}",
                )
                if st.button("Build Glossary", key=f"build_glossary_btn_{cp}"):
                    if not glossary_doc_names:
                        st.warning("Select at least one document.")
                    else:
                        parts = []
                        for name in glossary_doc_names:
                            d = next((doc for doc in proj["documents"] if doc["name"] == name), None)
                            if d:
                                parts.append(f"--- {name} ---\n{d['text']}")
                        with st.spinner("Extracting terms and definitions..."):
                            terms = generate_glossary("\n\n".join(parts))
                        if terms:
                            proj["glossary"] = terms
                            proj["glossary_terms_found"] = len(terms)

                terms = proj.get("glossary", [])
                if terms:
                    gl_df = pd.DataFrame(terms)
                    gl_df = gl_df.rename(columns={
                        "term": "Term", "definition": "Definition", "source_context": "Source Context",
                    })
                    gl_df = gl_df.sort_values("Term")
                    edited_gl_df = st.data_editor(gl_df, use_container_width=True, num_rows="dynamic", key=f"glossary_editor_{cp}")

                    gdl1, gdl2 = st.columns(2)
                    with gdl1:
                        st.download_button(
                            "Download as Excel (.xlsx)", build_xlsx_from_df("Glossary", edited_gl_df),
                            file_name="business_glossary.xlsx", mime=XLSX_MIME, key=f"gl_xlsx_{cp}",
                        )
                    with gdl2:
                        st.download_button(
                            "Download as Word (.docx)", build_docx_table_from_df("Business Glossary", edited_gl_df),
                            file_name="business_glossary.docx", mime=DOCX_MIME, key=f"gl_docx_{cp}",
                        )
                else:
                    st.info("Build the glossary to see terms here.")

    # --- Tab 1: Meeting Intelligence ---
    with tab1:
        section_header(
            "Meeting Intelligence & Actionizer",
            "Transform raw meeting transcripts into structured minutes, decisions, and action items.",
        )

        uploaded_transcript = st.file_uploader(
            "Upload Meeting Transcript (.txt or .docx):", type=['txt', 'docx'], key=f"transcript_uploader_{cp}",
        )

        if st.button("Process Transcript", key=f"process_transcript_btn_{cp}"):
            if uploaded_transcript is not None:
                with st.spinner("Extracting text..."):
                    text = extract_text_from_upload(uploaded_transcript)
                if not text.strip():
                    st.error("Couldn't extract any readable text from this transcript.")
                else:
                    with st.spinner("Extracting decisions, owners, and actions..."):
                        result = process_meeting(text)
                    if result:
                        proj["meeting_result"] = result
            else:
                st.warning("Please upload a transcript to process.")

        result = proj.get("meeting_result")
        if result:
            st.success("Meeting summary generated.")
            st.markdown("### Executive Summary")
            st.info(result.get("summary", ""))

            decisions = result.get("decisions", [])
            if decisions:
                st.markdown("### Key Decisions")
                for d in decisions:
                    st.write(f"- {d}")

            st.markdown("### Action Items Extracted")
            items = result.get("action_items", [])
            if items:
                action_df = pd.DataFrame(items)
                st.dataframe(action_df, use_container_width=True)
                dl1, dl2 = st.columns(2)
                with dl1:
                    st.download_button(
                        "Download as Excel (.xlsx)", build_xlsx_from_df("Action Items", action_df),
                        file_name="meeting_action_items.xlsx", mime=XLSX_MIME, key=f"meeting_xlsx_{cp}",
                    )
                with dl2:
                    st.download_button(
                        "Download as CSV (.csv)", action_df.to_csv(index=False),
                        file_name="meeting_action_items.csv", mime="text/csv", key=f"meeting_csv_{cp}",
                    )
            else:
                st.caption("No explicit action items were detected in this transcript.")

    # --- Tab 2: Elicitation Analysis ---
    with tab2:
        section_header(
            "Elicitation Analysis & Gap Detector",
            "Upload raw notes or transcripts. AI will structure needs and identify open questions.",
        )

        with st.expander("Prep for a Stakeholder Workshop", expanded=False):
            st.caption("Get a time-boxed agenda and targeted questions before you walk into the room.")
            focus_area = st.text_input(
                "What's this workshop about?",
                placeholder="e.g., Payments workflow requirements for the FinTech migration",
                key=f"workshop_focus_{cp}",
            )
            if st.button("Generate Workshop Prep", key=f"workshop_prep_btn_{cp}"):
                with st.spinner("Drafting agenda and questions..."):
                    prep = generate_workshop_prep(
                        proj.get("description", ""), focus_area, proj.get("extracted_text", ""),
                    )
                if prep:
                    proj["workshop_prep"] = prep

            prep = proj.get("workshop_prep")
            if prep:
                st.info(prep.get("objectives", ""))
                st.markdown("**Agenda**")
                agenda_df = pd.DataFrame(prep.get("agenda", []))
                if not agenda_df.empty:
                    agenda_df = agenda_df.rename(columns={
                        "topic": "Topic", "duration_minutes": "Minutes", "purpose": "Purpose",
                    })
                    st.dataframe(agenda_df, use_container_width=True, hide_index=True)
                st.markdown("**Questions to Ask**")
                for cat in sorted(set(q.get("category", "Other") for q in prep.get("questions", []))):
                    st.markdown(f"*{cat}*")
                    for q in prep.get("questions", []):
                        if q.get("category", "Other") == cat:
                            st.write(f"- {q.get('question', '')}")
                prep_md = f"# Workshop Prep — {focus_area or cp}\n\n{prep.get('objectives', '')}\n\n## Agenda\n"
                for a in prep.get("agenda", []):
                    prep_md += f"- **{a.get('topic', '')}** ({a.get('duration_minutes', '')} min) — {a.get('purpose', '')}\n"
                prep_md += "\n## Questions\n"
                for q in prep.get("questions", []):
                    prep_md += f"- [{q.get('category', '')}] {q.get('question', '')}\n"
                st.download_button(
                    "Download Workshop Prep (.md)", prep_md,
                    file_name="workshop_prep.md", mime="text/markdown", key=f"workshop_dl_{cp}",
                )

        st.markdown("---")

        repo_doc_names = [d["name"] for d in proj["documents"]]
        selected_repo_docs = []
        if repo_doc_names:
            selected_repo_docs = st.multiselect(
                "Include documents already in this project's repository:",
                repo_doc_names, default=repo_doc_names, key=f"elicit_repo_select_{cp}",
            )

        uploaded_file = st.file_uploader(
            "Upload a new Notes/Transcript or Document for this analysis:",
            type=['txt', 'pdf', 'docx', 'pptx', 'xlsx', 'csv'], key=f"gap_uploader_{cp}",
        )
        st.info(
            "Note on file intake: for proprietary formats (e.g., Apple Pages/Numbers, Visio, or "
            "live Google Docs/Sheets), please export to a universal format like .docx, .pdf, or "
            ".txt before uploading."
        )

        save_to_repo = False
        if uploaded_file is not None:
            save_to_repo = st.checkbox(
                "Also add this file to the project's document repository", value=True, key=f"gap_save_to_repo_{cp}",
            )

        notes = st.text_area(
            "Additional Notes / Context (optional):",
            placeholder=(
                "e.g., Focus on the payments workflow. The budget ceiling of $1M was confirmed "
                "by the sponsor on 6/10 — flag anything that conflicts with it."
            ),
            key=f"elicit_notes_{cp}", height=100,
        )
        if notes and notes.strip():
            st.caption("These notes will be factored into the analysis.")

        if st.button("Analyze for Gaps", key=f"analyze_gaps_btn_{cp}"):
            combined_parts = []
            for name in selected_repo_docs:
                doc = next((d for d in proj["documents"] if d["name"] == name), None)
                if doc:
                    combined_parts.append(f"--- Repository Document: {name} ---\n{doc['text']}")

            if uploaded_file is not None:
                with st.spinner("Extracting text from uploaded document..."):
                    new_text = extract_text_from_upload(uploaded_file)
                if new_text.strip():
                    combined_parts.append(f"--- Uploaded Document: {uploaded_file.name} ---\n{new_text}")
                    if save_to_repo:
                        add_doc_to_repo(proj, uploaded_file.name, new_text, uploaded_file.name.split(".")[-1].lower())

            combined_text = "\n\n".join(combined_parts)

            if not combined_text.strip():
                st.warning("Upload a document or select at least one repository document to analyze.")
            else:
                proj["extracted_text"] = combined_text
                proj["last_notes"] = notes or ""
                with st.spinner("Cross-referencing against requirements quality standards..."):
                    result = analyze_gaps(combined_text, notes=notes)
                if result:
                    proj["gap_analysis"] = result

        result = proj.get("gap_analysis")
        if result:
            n_open = len(result.get("open_questions", []))
            st.success(f"Analysis complete. Found {n_open} open item(s) for stakeholder follow-up.")
            if proj.get("last_notes"):
                st.caption(f"Notes accounted for: \"{proj['last_notes'][:200]}\"")
            if result.get("summary"):
                st.caption(result["summary"])
            st.metric(
                label="Requirements Risk Score",
                value=f"{result.get('risk_score', 0)}/100 ({result.get('risk_level', 'Unknown')})",
            )

            st.markdown("### Open Questions for Stakeholders")
            if n_open == 0:
                st.info("No significant gaps detected in this content.")
            for q in result.get("open_questions", []):
                st.warning(f"**{q.get('type', 'Issue')}:** {q.get('issue', '')}\n\n*Why it matters:* {q.get('why_it_matters', '')}")

        st.markdown("---")
        with st.expander("Prioritize Requirements (MoSCoW)", expanded=False):
            st.caption(
                "Scores the requirements in the content analyzed above as Must/Should/Could/Won't Have, "
                "with a rationale for each."
            )
            if st.button("Prioritize Requirements", key=f"prioritize_btn_{cp}"):
                source = proj.get("extracted_text", "")
                if not source.strip():
                    st.warning("Run 'Analyze for Gaps' above first (or upload/select documents) so there's content to prioritize.")
                else:
                    with st.spinner("Scoring requirements..."):
                        items = generate_prioritization(source)
                    if items:
                        proj["prioritization"] = items

            prioritized = proj.get("prioritization")
            if prioritized:
                pr_df = pd.DataFrame(prioritized)
                pr_df = pr_df.rename(columns={
                    "requirement": "Requirement", "moscow_category": "MoSCoW", "rationale": "Rationale",
                })
                category_order = {"Must Have": 0, "Should Have": 1, "Could Have": 2, "Won't Have": 3}
                pr_df["_sort"] = pr_df["MoSCoW"].map(category_order).fillna(4)
                pr_df = pr_df.sort_values("_sort").drop(columns="_sort")
                st.dataframe(pr_df, use_container_width=True, hide_index=True)
                dl1, dl2 = st.columns(2)
                with dl1:
                    st.download_button(
                        "Download as Excel (.xlsx)", build_xlsx_from_df("Prioritization", pr_df),
                        file_name="requirements_prioritization.xlsx", mime=XLSX_MIME, key=f"prio_xlsx_{cp}",
                    )
                with dl2:
                    st.download_button(
                        "Download as CSV (.csv)", pr_df.to_csv(index=False),
                        file_name="requirements_prioritization.csv", mime="text/csv", key=f"prio_csv_{cp}",
                    )

    # --- Tab 3: Documentation Generator ---
    with tab3:
        section_header(
            "Documentation Generator",
            "Generate a real first-draft document from your requirements notes, exportable to Word or Excel.",
        )

        doc_type = st.selectbox("Select Document Type to Draft", list(DOC_TYPE_CODES.keys()), key=f"doc_type_select_{cp}")

        repo_doc_names = [d["name"] for d in proj["documents"]]
        selected_repo_docs = []
        if repo_doc_names:
            selected_repo_docs = st.multiselect(
                "Include documents from this project's repository:",
                repo_doc_names, default=repo_doc_names, key=f"doc_repo_select_{cp}",
            )

        repo_text_parts = []
        for name in selected_repo_docs:
            doc = next((d for d in proj["documents"] if d["name"] == name), None)
            if doc:
                repo_text_parts.append(f"--- {name} ---\n{doc['text']}")
        repo_combined = "\n\n".join(repo_text_parts)

        context_text = st.text_area(
            "Additional source content / notes (combined with the repository documents selected above):",
            value="" if repo_doc_names else proj.get("extracted_text", "")[:2000],
            height=120, key=f"doc_context_{cp}",
        )
        full_context = f"{repo_combined}\n\n{context_text}".strip() if repo_combined else context_text

        user_suggestion = st.text_area(
            "Provide specific instructions or focus areas:",
            "e.g., Ensure the regulatory compliance section is highly detailed.",
            key=f"doc_suggestion_{cp}",
        )

        if st.button(f"Generate Draft {doc_type}", key=f"generate_doc_btn_{cp}"):
            if doc_type == "Data Dictionary":
                with st.spinner("Drafting data dictionary..."):
                    rows = generate_data_dictionary(full_context, user_suggestion)
                if rows:
                    proj["documents_drafted"] = proj.get("documents_drafted", 0) + 1
                    proj["last_doc_draft"] = {"kind": "data_dictionary", "rows": rows}
                    proj["last_doc_type"] = doc_type
            elif doc_type == "As-Is / To-Be Process Document":
                with st.spinner("Mapping As-Is and To-Be process steps..."):
                    rows = generate_asis_tobe(full_context, user_suggestion)
                if rows:
                    proj["documents_drafted"] = proj.get("documents_drafted", 0) + 1
                    proj["last_doc_draft"] = {"kind": "asis_tobe", "rows": rows}
                    proj["last_doc_type"] = doc_type
            else:
                with st.spinner(f"Drafting {doc_type}..."):
                    draft = generate_document(doc_type, full_context, user_suggestion)
                if draft:
                    proj["documents_drafted"] = proj.get("documents_drafted", 0) + 1
                    proj["last_doc_draft"] = {"kind": "markdown", "text": draft}
                    proj["last_doc_type"] = doc_type

        saved = proj.get("last_doc_draft")
        if saved:
            shown_type = proj.get("last_doc_type", doc_type)
            code = DOC_TYPE_CODES.get(shown_type, "Document")
            st.success(f"Draft of {shown_type} generated.")

            if saved["kind"] == "markdown":
                st.markdown(saved["text"])
                dl1, dl2 = st.columns(2)
                with dl1:
                    st.download_button(
                        "Download as Word (.docx)", build_docx_from_markdown(shown_type, saved["text"]),
                        file_name=f"{code}.docx", mime=DOCX_MIME, key=f"doc_md_docx_{cp}",
                    )
                with dl2:
                    st.download_button(
                        "Download as Markdown (.md)", saved["text"], file_name=f"{code}.md", mime="text/markdown",
                        key=f"doc_md_md_{cp}",
                    )

            elif saved["kind"] == "data_dictionary":
                df = pd.DataFrame(saved["rows"])
                df = df.rename(columns={
                    "field_name": "Field Name", "data_type": "Data Type", "description": "Description",
                    "source_system": "Source System", "validation_rules": "Validation Rules",
                })
                df = df.reindex(columns=["Field Name", "Data Type", "Description", "Source System", "Validation Rules"], fill_value="")
                st.dataframe(df, use_container_width=True)
                dl1, dl2 = st.columns(2)
                with dl1:
                    st.download_button(
                        "Download as Excel (.xlsx)", build_xlsx_from_df("Data Dictionary", df),
                        file_name=f"{code}.xlsx", mime=XLSX_MIME, key=f"doc_dd_xlsx_{cp}",
                    )
                with dl2:
                    st.download_button(
                        "Download as Word (.docx)", build_docx_table_from_df("Data Dictionary", df),
                        file_name=f"{code}.docx", mime=DOCX_MIME, key=f"doc_dd_docx_{cp}",
                    )

            elif saved["kind"] == "asis_tobe":
                df = pd.DataFrame(saved["rows"])
                df = df.rename(columns={
                    "step_id": "Step ID", "process_step": "Process Step", "as_is_description": "As-Is",
                    "to_be_description": "To-Be", "gap_or_change": "Gap / Change Needed",
                    "shape_type": "Shape Type", "next_step_id": "Next Step ID",
                })
                col_order = ["Step ID", "Process Step", "As-Is", "To-Be", "Gap / Change Needed", "Shape Type", "Next Step ID"]
                df = df.reindex(columns=col_order, fill_value="")
                st.dataframe(df, use_container_width=True)

                dl1, dl2, dl3 = st.columns(3)
                with dl1:
                    st.download_button(
                        "Download as Excel (.xlsx)", build_xlsx_from_df("As-Is To-Be", df),
                        file_name=f"{code}.xlsx", mime=XLSX_MIME, key=f"doc_at_xlsx_{cp}",
                    )
                with dl2:
                    st.download_button(
                        "Download as Word (.docx)", build_docx_table_from_df("As-Is / To-Be Process Document", df),
                        file_name=f"{code}.docx", mime=DOCX_MIME, key=f"doc_at_docx_{cp}",
                    )
                with dl3:
                    st.download_button(
                        "Download Visio Process Map (.xlsx)", build_visio_dataviz_xlsx(saved["rows"]),
                        file_name=f"{code}_VisioDataVisualizer.xlsx", mime=XLSX_MIME, key=f"doc_at_visio_{cp}",
                    )
                st.caption(
                    "The Visio Process Map file is formatted for Visio's built-in Data Visualizer "
                    "feature: in Visio, start a Data Visualizer Basic Flowchart template and import "
                    "this file to auto-generate an editable diagram. (Native .vsdx files can't be "
                    "reliably authored from scratch without Visio itself or a paid SDK — this Excel-based "
                    "import is the documented, reliable path to a real Visio diagram.)"
                )

    # --- Tab 4: Agile Story Creator ---
    with tab4:
        section_header(
            "Agile Story & Backlog Creator",
            "Convert validated requirements into ready-to-import User Stories and Gherkin Acceptance Criteria.",
        )

        repo_doc_names = [d["name"] for d in proj["documents"]]
        selected_repo_docs = []
        if repo_doc_names:
            selected_repo_docs = st.multiselect(
                "Include documents from this project's repository:",
                repo_doc_names, default=repo_doc_names, key=f"story_repo_select_{cp}",
            )

        repo_text_parts = []
        for name in selected_repo_docs:
            doc = next((d for d in proj["documents"] if d["name"] == name), None)
            if doc:
                repo_text_parts.append(f"--- {name} ---\n{doc['text']}")
        repo_combined = "\n\n".join(repo_text_parts)

        notes_text = st.text_area(
            "Additional requirements / notes (combined with the repository documents selected above):",
            value="" if repo_doc_names else proj.get("extracted_text", "")[:2000],
            height=120, key=f"story_source_{cp}",
        )
        source_text = f"{repo_combined}\n\n{notes_text}".strip() if repo_combined else notes_text

        if st.button("Generate User Stories & Acceptance Criteria", key=f"generate_stories_btn_{cp}"):
            if not source_text.strip():
                st.warning("Add some requirements text, or select at least one repository document, first.")
            else:
                with st.spinner("Drafting user stories and acceptance criteria..."):
                    stories = generate_stories(source_text)
                if stories:
                    proj["stories"] = stories
                    proj["stories_drafted"] = proj.get("stories_drafted", 0) + len(stories)

        stories = proj.get("stories", [])
        if stories:
            st.markdown("### User Story Drafts")
            story_df = pd.DataFrame(stories)
            edited_df = st.data_editor(story_df, use_container_width=True, num_rows="dynamic", key=f"story_editor_{cp}")

            dl1, dl2 = st.columns(2)
            with dl1:
                st.download_button(
                    "Download as Excel (.xlsx)", build_xlsx_from_df("Backlog", edited_df),
                    file_name="backlog_stories.xlsx", mime=XLSX_MIME, key=f"story_xlsx_{cp}",
                )
            with dl2:
                st.download_button(
                    "Download as CSV — Jira/Azure DevOps import format (.csv)", edited_df.to_csv(index=False),
                    file_name="backlog_stories.csv", mime="text/csv", key=f"story_csv_{cp}",
                )

            st.markdown("---")
            with st.expander("Generate Test Cases from These Stories", expanded=False):
                st.caption("Derives test cases straight from the acceptance criteria above — no re-typing.")
                if st.button("Generate Test Cases", key=f"generate_tc_btn_{cp}"):
                    with st.spinner("Writing test cases..."):
                        test_cases = generate_test_cases(edited_df.to_dict("records"))
                    if test_cases:
                        proj["test_cases"] = test_cases
                        proj["test_cases_drafted"] = proj.get("test_cases_drafted", 0) + len(test_cases)

                test_cases = proj.get("test_cases", [])
                if test_cases:
                    tc_df = pd.DataFrame(test_cases)
                    tc_df = tc_df.rename(columns={
                        "test_id": "Test ID", "related_story": "Related Story", "scenario": "Scenario",
                        "preconditions": "Preconditions", "steps": "Steps",
                        "expected_result": "Expected Result", "priority": "Priority",
                    })
                    col_order = ["Test ID", "Related Story", "Scenario", "Preconditions", "Steps", "Expected Result", "Priority"]
                    tc_df = tc_df.reindex(columns=col_order, fill_value="")
                    edited_tc_df = st.data_editor(tc_df, use_container_width=True, num_rows="dynamic", key=f"tc_editor_{cp}")

                    tdl1, tdl2 = st.columns(2)
                    with tdl1:
                        st.download_button(
                            "Download as Excel (.xlsx)", build_xlsx_from_df("Test Cases", edited_tc_df),
                            file_name="test_cases.xlsx", mime=XLSX_MIME, key=f"tc_xlsx_{cp}",
                        )
                    with tdl2:
                        st.download_button(
                            "Download as CSV (.csv)", edited_tc_df.to_csv(index=False),
                            file_name="test_cases.csv", mime="text/csv", key=f"tc_csv_{cp}",
                        )
                else:
                    st.info("Generate test cases to see them here.")
        else:
            st.info("Generate stories to see them here.")

    # --- Tab 5: Traceability & Change Impact ---
    with tab5:
        section_header(
            "Requirements Traceability Matrix",
            "Auto-links Requirement → User Story → Test Case → Priority from what you've already generated. Edit freely.",
        )

        if not proj.get("stories"):
            st.info("Generate user stories in the Agile Story & Backlog Creator tab first — the RTM builds from those.")
        else:
            if st.button("Build / Refresh RTM", key=f"build_rtm_btn_{cp}"):
                proj["rtm_rows"] = build_rtm_rows(proj)

            rtm_rows = proj.get("rtm_rows", [])
            if not rtm_rows:
                st.info("Click 'Build / Refresh RTM' to auto-populate from your stories, test cases, and prioritization.")
            else:
                rtm_df = pd.DataFrame(rtm_rows)
                rtm_df = rtm_df.rename(columns={
                    "requirement": "Requirement", "user_story": "User Story", "test_cases": "Test Case(s)",
                    "priority": "Priority", "status": "Status",
                })
                col_order = ["Requirement", "User Story", "Test Case(s)", "Priority", "Status"]
                rtm_df = rtm_df.reindex(columns=col_order, fill_value="")
                st.caption(
                    "Auto-matching is by requirement label, so it may miss or misalign a few links — "
                    "add, remove, or fix rows directly in the table below before exporting."
                )
                edited_rtm_df = st.data_editor(rtm_df, use_container_width=True, num_rows="dynamic", key=f"rtm_editor_{cp}")

                dl1, dl2 = st.columns(2)
                with dl1:
                    st.download_button(
                        "Download as Excel (.xlsx)", build_xlsx_from_df("RTM", edited_rtm_df),
                        file_name="requirements_traceability_matrix.xlsx", mime=XLSX_MIME, key=f"rtm_xlsx_{cp}",
                    )
                with dl2:
                    st.download_button(
                        "Download as Word (.docx)", build_docx_table_from_df("Requirements Traceability Matrix", edited_rtm_df),
                        file_name="requirements_traceability_matrix.docx", mime=DOCX_MIME, key=f"rtm_docx_{cp}",
                    )

        st.markdown("---")
        section_header(
            "Change Request Impact Analyzer",
            "Paste an incoming change request. AI checks it against this project's requirements/stories and flags what it touches.",
        )

        change_request_text = st.text_area(
            "Describe the change request:",
            placeholder=(
                "e.g., Sponsor now wants multi-currency support added to the payments workflow "
                "before launch, on top of the original single-currency scope."
            ),
            key=f"change_request_{cp}", height=100,
        )

        if st.button("Analyze Change Impact", key=f"analyze_change_btn_{cp}"):
            if not change_request_text.strip():
                st.warning("Describe the change request first.")
            else:
                existing_parts = []
                if proj.get("extracted_text"):
                    existing_parts.append(f"--- Elicitation Source Content ---\n{proj['extracted_text']}")
                if proj.get("stories"):
                    story_lines = "\n".join(
                        f"- {s.get('requirement', '')}: {s.get('user_story', '')}" for s in proj["stories"]
                    )
                    existing_parts.append(f"--- Generated User Stories ---\n{story_lines}")
                existing_context = "\n\n".join(existing_parts)

                with st.spinner("Assessing change impact..."):
                    impact = generate_change_impact(change_request_text, existing_context)
                if impact:
                    history = proj.get("change_impact_history", [])
                    history.append({"request_text": change_request_text, "result": impact})
                    proj["change_impact_history"] = history

        change_history = proj.get("change_impact_history", [])
        if change_history:
            latest = change_history[-1]
            impact = latest["result"]
            level_color = {"Low": "info", "Medium": "warning", "High": "warning", "Critical": "error"}.get(
                impact.get("impact_level", "Medium"), "info"
            )
            st.markdown(f"**Change request:** {latest['request_text']}")
            getattr(st, level_color)(f"**Impact Level: {impact.get('impact_level', 'Unknown')}** — {impact.get('summary', '')}")

            affected = impact.get("affected_requirements", [])
            if affected:
                st.markdown("**Affected Requirements/Stories**")
                for a in affected:
                    st.write(f"- **{a.get('requirement', '')}:** {a.get('impact_description', '')}")
            else:
                st.caption("No existing requirements/stories were flagged as affected.")

            actions = impact.get("recommended_actions", [])
            if actions:
                st.markdown("**Recommended Next Actions**")
                for act in actions:
                    st.write(f"- {act}")

            if len(change_history) > 1:
                with st.expander(f"Previous change requests analyzed ({len(change_history) - 1})"):
                    for prev in reversed(change_history[:-1]):
                        st.markdown(f"**{prev['request_text']}**")
                        st.caption(f"{prev['result'].get('impact_level', '')} — {prev['result'].get('summary', '')}")
                        st.markdown("---")


def pm_module():
    section_header(
        "Project Managers: Predictive Risk & Health (Placeholder)",
        "View predictive metrics, resource optimization, and automated status reports.",
    )
    st.selectbox("Select Project to View", list(st.session_state["projects"].keys()))
    st.info("PM features (Project Health Forecaster, Constraint Solver) haven't been built yet — this module is still a placeholder.")


def pgm_module():
    section_header(
        "Program Managers: Portfolio Optimization (Placeholder)",
        "Analyze cross-project dependencies, resource contention, and benefit realization.",
    )
    st.warning("PgM features (Interdependency Mapper, Benefit Realization Tracker) haven't been built yet — this module is still a placeholder.")


if __name__ == "__main__":
    # --- Main App Navigation ---
    import auth  # local import — see note near check_access()

    # Theme goes in before the auth gate — otherwise the login/signup screen
    # renders completely unstyled, since check_access() st.stop()s before any
    # of this CSS would have been injected.
    inject_theme()

    if not check_access():
        st.stop()

    render_masthead()
    init_projects()

    # Apply any pending project switch (e.g. from creating a new project) BEFORE the
    # sidebar selectbox below is instantiated — Streamlit won't allow setting a widget's
    # session_state value after that widget has already rendered in the same run.
    if "pending_project_switch" in st.session_state:
        _target = st.session_state.pop("pending_project_switch")
        if _target in st.session_state["projects"]:
            st.session_state["current_project"] = _target

    st.sidebar.markdown(
        f"<div style='font-size:1.6rem; font-weight:800; color:{ACCENT}; letter-spacing:0.5px;'>ScopeForge</div>"
        f"<div style='font-size:0.8rem; color:{SIDEBAR_MUTED}; margin-bottom:0.8rem;'>Consulting Accelerator</div>",
        unsafe_allow_html=True,
    )

    auth.render_logout_control()
    st.sidebar.markdown("---")

    st.sidebar.subheader("Active Project")
    project_names = list(st.session_state["projects"].keys())
    st.sidebar.selectbox("Select Project", project_names, key="current_project")

    st.sidebar.markdown("---")
    st.sidebar.subheader("Modules")
    role = st.sidebar.radio("Select Your Role", ["Business Analyst (BA)", "Project Manager (PM)", "Program Manager (PgM)"])

    st.sidebar.markdown("---")
    st.sidebar.selectbox(
        "AI Model", [DEFAULT_MODEL, FAST_MODEL], index=0, key="model",
        format_func=lambda m: MODEL_LABELS.get(m, m),
        help="Sonnet = best quality for analysis/drafting. Haiku = faster and cheaper, good for quick checks.",
    )
    if not get_api_key():
        st.sidebar.error("No ANTHROPIC_API_KEY found in Secrets.")
    else:
        st.sidebar.success("API key loaded.")

    st.sidebar.markdown("---")

    # --- Display Selected Module ---
    if role == "Business Analyst (BA)":
        ba_module()
    elif role == "Project Manager (PM)":
        pm_module()
    elif role == "Program Manager (PgM)":
        pgm_module()

    st.divider()
    section_header("ScopeBot (AI Assistant)", "Ask about requirements, JIRA sync, or BA best practices.")

    if "chat_history" not in st.session_state:
        st.session_state["chat_history"] = []

    for msg in st.session_state["chat_history"]:
        st.chat_message(msg["role"]).write(msg["content"])

    user_query = st.chat_input("Ask ScopeBot a question about requirements, JIRA sync, or best practices...")

    if user_query:
        st.session_state["chat_history"].append({"role": "user", "content": user_query})
        st.chat_message("user").write(user_query)

        with st.chat_message("assistant"):
            with st.spinner("Thinking..."):
                reply = chat_with_bot(st.session_state["chat_history"])
            if reply:
                st.write(reply)
            else:
                reply = "Sorry, I couldn't process that — please try again."
                st.write(reply)

        st.session_state["chat_history"].append({"role": "assistant", "content": reply})
